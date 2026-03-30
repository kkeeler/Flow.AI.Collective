"""
Assemble a PPTX from captured screenshots.

Supports static slides (1 image) and animated slides (layered images with
click-to-advance XML). Adapt SLIDE_DEFS to your project.

    python build_pptx.py
"""

import io
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Emu

SCREENSHOT_DIR = "screenshots"
OUTPUT_FILE = "output.pptx"

SW = Inches(13.333)
SH = Inches(7.5)

# Define slides in order. Each entry is either:
#   ("static", "screenshot_name.png")
#   ("animated", "base_name", step_count)  — expects base_name_step0.png .. base_name_stepN.png
SLIDE_DEFS = [
    ("static", "s1.png"),
    ("animated", "s2", 5),
    ("static", "s3.png"),
]


def add_static_slide(prs, image_path):
    """Add a slide with a single full-bleed image."""
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)
    slide.shapes.add_picture(str(image_path), 0, 0, SW, SH)
    return slide


def add_animated_slide(prs, base_name, steps):
    """Add a slide with layered images and click-to-advance animation."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    img_dir = Path(SCREENSHOT_DIR)

    shape_ids = []
    for i in range(steps + 1):
        img = img_dir / f"{base_name}_step{i}.png"
        pic = slide.shapes.add_picture(str(img), 0, 0, SW, SH)
        shape_ids.append(pic.shape_id)

    _inject_animation_xml(slide, shape_ids[1:])  # animate all except base
    return slide


def _inject_animation_xml(slide, spids):
    """Inject click-to-advance animation XML for the given shape IDs."""
    NSMAP = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    ns = NSMAP["p"]

    ctn_id = [1]

    def next_id():
        val = ctn_id[0]
        ctn_id[0] += 1
        return str(val)

    def make_click_par(spid):
        """Build one <p:par> block for a single click step."""
        par = etree.SubElement(child_tn_lst, f"{{{ns}}}par")
        ctn = etree.SubElement(par, f"{{{ns}}}cTn", id=next_id(), fill="hold")
        st = etree.SubElement(ctn, f"{{{ns}}}stCondLst")
        etree.SubElement(st, f"{{{ns}}}cond", delay="indefinite")
        inner_child = etree.SubElement(ctn, f"{{{ns}}}childTnLst")

        par2 = etree.SubElement(inner_child, f"{{{ns}}}par")
        ctn2 = etree.SubElement(par2, f"{{{ns}}}cTn", id=next_id(), fill="hold")
        st2 = etree.SubElement(ctn2, f"{{{ns}}}stCondLst")
        etree.SubElement(st2, f"{{{ns}}}cond", delay="0")
        inner2 = etree.SubElement(ctn2, f"{{{ns}}}childTnLst")

        par3 = etree.SubElement(inner2, f"{{{ns}}}par")
        ctn3 = etree.SubElement(par3, f"{{{ns}}}cTn", id=next_id(),
                                presetID="10", presetClass="entr",
                                presetSubtype="0", fill="hold", nodeType="clickEffect")
        st3 = etree.SubElement(ctn3, f"{{{ns}}}stCondLst")
        etree.SubElement(st3, f"{{{ns}}}cond", delay="0")
        anim_children = etree.SubElement(ctn3, f"{{{ns}}}childTnLst")

        # <p:set> — make visible
        p_set = etree.SubElement(anim_children, f"{{{ns}}}set")
        cbhvr = etree.SubElement(p_set, f"{{{ns}}}cBhvr")
        set_ctn = etree.SubElement(cbhvr, f"{{{ns}}}cTn", id=next_id(), dur="1", fill="hold")
        set_st = etree.SubElement(set_ctn, f"{{{ns}}}stCondLst")
        etree.SubElement(set_st, f"{{{ns}}}cond", delay="0")
        tgt = etree.SubElement(cbhvr, f"{{{ns}}}tgtEl")
        etree.SubElement(tgt, f"{{{ns}}}spTgt", spid=str(spid))
        attr_lst = etree.SubElement(cbhvr, f"{{{ns}}}attrNameLst")
        attr_name = etree.SubElement(attr_lst, f"{{{ns}}}attrName")
        attr_name.text = "style.visibility"
        to_el = etree.SubElement(p_set, f"{{{ns}}}to")
        etree.SubElement(to_el, f"{{{ns}}}strVal", val="visible")

        # <p:animEffect> — fade in
        anim_eff = etree.SubElement(anim_children, f"{{{ns}}}animEffect",
                                    transition="in", filter="fade")
        cbhvr2 = etree.SubElement(anim_eff, f"{{{ns}}}cBhvr")
        etree.SubElement(cbhvr2, f"{{{ns}}}cTn", id=next_id(), dur="500")
        tgt2 = etree.SubElement(cbhvr2, f"{{{ns}}}tgtEl")
        etree.SubElement(tgt2, f"{{{ns}}}spTgt", spid=str(spid))

    # Build top-level timing tree
    timing = etree.SubElement(slide._element, f"{{{ns}}}timing")
    tn_lst = etree.SubElement(timing, f"{{{ns}}}tnLst")
    root_par = etree.SubElement(tn_lst, f"{{{ns}}}par")
    root_ctn = etree.SubElement(root_par, f"{{{ns}}}cTn", id=next_id(),
                                dur="indefinite", restart="never", nodeType="tmRoot")
    root_children = etree.SubElement(root_ctn, f"{{{ns}}}childTnLst")

    seq = etree.SubElement(root_children, f"{{{ns}}}seq", concurrent="1", nextAc="seek")
    seq_ctn = etree.SubElement(seq, f"{{{ns}}}cTn", id=next_id(),
                               dur="indefinite", nodeType="mainSeq")
    child_tn_lst = etree.SubElement(seq_ctn, f"{{{ns}}}childTnLst")

    for spid in spids:
        make_click_par(spid)

    # Navigation conditions
    prev_cond_lst = etree.SubElement(seq, f"{{{ns}}}prevCondLst")
    prev_cond = etree.SubElement(prev_cond_lst, f"{{{ns}}}cond", evt="onPrev", delay="0")
    prev_tgt = etree.SubElement(prev_cond, f"{{{ns}}}tgtEl")
    etree.SubElement(prev_tgt, f"{{{ns}}}sldTgt")

    next_cond_lst = etree.SubElement(seq, f"{{{ns}}}nextCondLst")
    next_cond = etree.SubElement(next_cond_lst, f"{{{ns}}}cond", evt="onNext", delay="0")
    next_tgt = etree.SubElement(next_cond, f"{{{ns}}}tgtEl")
    etree.SubElement(next_tgt, f"{{{ns}}}sldTgt")


def set_speaker_notes(slide, text):
    """Set speaker notes on a slide, creating the notes slide if needed."""
    if not slide.has_notes_slide:
        slide.notes_slide  # force creation
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    for defn in SLIDE_DEFS:
        if defn[0] == "static":
            slide = add_static_slide(prs, Path(SCREENSHOT_DIR) / defn[1])
            print(f"  Added static slide: {defn[1]}")
        elif defn[0] == "animated":
            slide = add_animated_slide(prs, defn[1], defn[2])
            print(f"  Added animated slide: {defn[1]} ({defn[2]} clicks)")

    prs.save(OUTPUT_FILE)
    print(f"\nSaved to {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
